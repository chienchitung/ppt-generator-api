import json
import argparse
import sys
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_THEME_COLOR
import locale
import requests
from io import BytesIO
from datetime import datetime
import warnings
from PIL import Image
import tempfile
import os

# Suppress InsecureRequestWarning
warnings.filterwarnings("ignore", category=requests.packages.urllib3.exceptions.InsecureRequestWarning)

# Set locale to handle Chinese characters
try:
    locale.setlocale(locale.LC_ALL, 'zh_TW.UTF-8')
except locale.Error:
    try:
        locale.setlocale(locale.LC_ALL, 'zh_TW')
    except locale.Error:
        try:
            locale.setlocale(locale.LC_ALL, '')  # Use system default locale
        except locale.Error:
            logger.warning("Could not set Chinese locale. Some characters might not display correctly.")

# Setup logging with UTF-8 encoding
logging.basicConfig(level=logging.INFO, encoding='utf-8')
logger = logging.getLogger(__name__)

def convert_webp_to_png(webp_data):
    """Convert WebP image to PNG format"""
    try:
        image = Image.open(webp_data)
        logger.info(f"Image opened successfully, format: {image.format}, mode: {image.mode}, size: {image.size}")
        
        # Convert to RGB mode if needed
        if image.mode in ('RGBA', 'LA') or (image.mode == 'P' and 'transparency' in image.info):
            logger.info("Converting image with alpha channel...")
            # Create a white background
            background = Image.new('RGBA', image.size, (255, 255, 255))
            # Paste the image using alpha channel as mask
            if image.mode == 'RGBA':
                background.paste(image, mask=image.split()[3])
            else:
                background.paste(image)
            # Convert to RGB
            image = background.convert('RGB')
        else:
            logger.info(f"Converting image without alpha channel from mode: {image.mode}")
            image = image.convert('RGB')
        
        # Save as PNG
        output = BytesIO()
        image.save(output, format="PNG")
        output.seek(0)
        logger.info("Image successfully converted to PNG")
        return output
    except Exception as e:
        logger.error(f"Error converting WebP to PNG: {str(e)}")
        # Try a different approach as fallback
        try:
            logger.info("Trying alternative conversion method...")
            image = Image.open(webp_data)
            output = BytesIO()
            # Try saving as JPEG if PNG fails
            image.convert('RGB').save(output, format="JPEG")
            output.seek(0)
            logger.info("Image successfully converted to JPEG as fallback")
            return output
        except Exception as e2:
            logger.error(f"Error in fallback conversion: {str(e2)}")
            return None

def download_image(url):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, verify=False, timeout=10)
        response.raise_for_status()
        image_data = BytesIO(response.content)
        
        # Check if it's a WebP image and convert if needed
        if url.lower().endswith('.webp') or '.webp' in url.lower():
            logger.info(f"Converting WebP image from {url} to PNG")
            return convert_webp_to_png(image_data)
        
        return image_data
    except requests.exceptions.SSLError as e:
        logger.error(f"SSL Error downloading image from {url}: {str(e)}")
        try:
            # Retry without SSL verification
            response = requests.get(url, headers=headers, verify=False, timeout=10)
            response.raise_for_status()
            image_data = BytesIO(response.content)
            
            # Check if it's a WebP image and convert if needed
            if url.lower().endswith('.webp') or '.webp' in url.lower():
                logger.info(f"Converting WebP image from {url} to PNG")
                return convert_webp_to_png(image_data)
            
            return image_data
        except Exception as e:
            logger.error(f"Error retrying download without SSL verification from {url}: {str(e)}")
            return None
    except requests.exceptions.RequestException as e:
        logger.error(f"Error downloading image from {url}: {str(e)}")
        return None
    except Exception as e:
        logger.error(f"Unexpected error downloading image from {url}: {str(e)}")
        return None

def set_slide_size_to_16x9(prs):
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    
    # Add background shape
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    background.line.fill.background()

    # Add decorative elements
    left_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(1), prs.slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RGBColor(30, 144, 255)
    left_bar.line.fill.background()

    # Add title
    title_box = shapes.add_textbox(
        Inches(2), Inches(3), Inches(12), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    title_para.alignment = PP_ALIGN.LEFT

    # Add subtitle (date)
    subtitle_box = shapes.add_textbox(
        Inches(2), Inches(4.5), Inches(12), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.LEFT

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    shapes = slide.shapes
    
    # Add gradient background
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)
    background.line.fill.background()

    # Add title
    title_box = shapes.add_textbox(
        Inches(2), Inches(3), Inches(12), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add decorative line
    line = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2),
        Inches(4.5),
        Inches(12),
        Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(30, 144, 255)
    line.line.fill.background()

def download_and_convert_image(url, app_name):
    """Download image and convert to PNG if needed, saving it to a temp file to ensure reliability"""
    try:
        logger.info(f"Downloading and processing image for {app_name} from: {url}")
        
        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        temp_input_path = os.path.join(temp_dir, f"{app_name}_original")
        temp_output_path = os.path.join(temp_dir, f"{app_name}_converted.png")
        
        # Download image
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, verify=False, timeout=10)
        response.raise_for_status()
        
        # Save original image
        with open(temp_input_path, 'wb') as f:
            f.write(response.content)
        
        # Open and convert
        with Image.open(temp_input_path) as img:
            logger.info(f"Image format: {img.format}, mode: {img.mode}, size: {img.size}")
            
            # Convert to RGB
            if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                background = Image.new('RGBA', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[3])
                else:
                    background.paste(img)
                img = background.convert('RGB')
            else:
                img = img.convert('RGB')
                
            # Save as PNG
            img.save(temp_output_path, 'PNG')
            logger.info(f"Image saved to: {temp_output_path}")
            
            # Return file path for adding to slide
            return temp_output_path
    except Exception as e:
        logger.error(f"Error in download_and_convert_image: {str(e)}")
        return None

def add_content_slide(prs, title, content, app_icon_url=None, app_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # Add subtle background
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Add title with app icon if available
    title_left = Inches(1)  # 標題預設位置
    title_width = Inches(13)  # 標題預設寬度
    
    if app_icon_url and app_name:
        try:
            # Use the local file conversion approach
            local_image_path = download_and_convert_image(app_icon_url, app_name)
            
            if local_image_path and os.path.exists(local_image_path):
                logger.info(f"Adding image from local path: {local_image_path}")
                # 將圖示放在右上角 (右邊減去圖片寬度減去間隔)
                icon = shapes.add_picture(
                    local_image_path,
                    prs.slide_width - Inches(1.5),  # 右邊位置
                    Inches(0.5),  # 上方位置
                    height=Inches(1)  # 高度
                )
                # 保持標題在左側，但縮短標題寬度避免與圖示重疊
                title_width = Inches(11)
                logger.info("Image successfully added to slide from local file at right corner")
                
                # Clean up temp files
                try:
                    os.remove(local_image_path)
                    # Remove the original image file as well
                    original_path = os.path.join(os.path.dirname(local_image_path), f"{app_name}_original")
                    if os.path.exists(original_path):
                        os.remove(original_path)
                    # Remove the temp directory
                    temp_dir = os.path.dirname(local_image_path)
                    if os.path.exists(temp_dir):
                        os.rmdir(temp_dir)
                except Exception as cleanup_error:
                    logger.warning(f"Error cleaning up temp files: {str(cleanup_error)}")
            else:
                logger.warning("Local image file not created, skipping add_picture")
        except Exception as e:
            logger.error(f"Error adding app icon from local file: {str(e)}")
    
    # Add title (不再移動標題位置，只調整寬度)
    title_box = shapes.add_textbox(
        title_left, Inches(0.5),
        title_width, Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    
    content_left = Inches(1)
    content_top = Inches(1.8)
    content_width = Inches(14)
    content_height = Inches(6)
    
    content_box = shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for key, value in content.items():
        p = tf.add_paragraph()
        p.text = f"{key}:"
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(60, 60, 60)
        
        if isinstance(value, list):
            for item in value:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.level = 1
                p.font.color.rgb = RGBColor(80, 80, 80)
        elif isinstance(value, dict):
            for sub_key, sub_value in value.items():
                p = tf.add_paragraph()
                p.text = f"• {sub_key}: {sub_value}"
                p.font.size = Pt(18)
                p.level = 1
                p.font.color.rgb = RGBColor(80, 80, 80)
        else:
            p = tf.add_paragraph()
            p.text = f"• {value}"
            p.font.size = Pt(18)
            p.level = 1
            p.font.color.rgb = RGBColor(80, 80, 80)

def add_comparison_slide(prs, title, apps_data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # Add background
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    title_shape = shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 144, 255)
    
    # Calculate layout for comparison boxes
    box_width = Inches(4.5)
    box_height = Inches(6)
    gap = Inches(0.5)
    start_left = (prs.slide_width - (box_width * len(apps_data) + gap * (len(apps_data) - 1))) / 2
    top = Inches(2)
    
    for i, app in enumerate(apps_data):
        left = start_left + (box_width + gap) * i
        box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(200, 200, 200)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = app["name"]
        p.font.bold = True
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        
        # Add content specific to comparison type
        add_comparison_content(tf, app)

def add_comparison_content(text_frame, app_data):
    # Add specific comparison content based on the data
    pass  # Implementation will be added based on specific comparison needs

def add_chapter_slide(prs, chapter_data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # Add background
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    # Add title
    title_box = shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = chapter_data["title"]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    
    # Add content sections
    sections = [
        ("關鍵發現", chapter_data["keyFindings"]),
        ("關鍵數據", chapter_data["dataSupport"]),
        ("建議", chapter_data["recommendations"])
    ]
    
    current_top = Inches(2)
    for section_title, section_content in sections:
        # Add section title
        section_title_box = shapes.add_textbox(
            Inches(1), current_top,
            Inches(14), Inches(0.5)
        )
        section_title_frame = section_title_box.text_frame
        section_title_para = section_title_frame.add_paragraph()
        section_title_para.text = section_title
        section_title_para.font.size = Pt(24)
        section_title_para.font.bold = True
        section_title_para.font.color.rgb = RGBColor(60, 60, 60)
        
        # Add section content
        content_box = shapes.add_textbox(
            Inches(1), current_top + Inches(0.7),
            Inches(14), Inches(1.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for item in section_content:
            p = content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(80, 80, 80)
        
        current_top += Inches(2.5)

def add_ending_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    
    # Add background shape
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background
    background.line.fill.background()

    # Add decorative elements
    left_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(1), prs.slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = RGBColor(30, 144, 255)
    left_bar.line.fill.background()

    # Add title
    title_box = shapes.add_textbox(
        Inches(2), Inches(3), Inches(12), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "謝謝聆聽"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    title_para.alignment = PP_ALIGN.LEFT

    # Add subtitle
    subtitle_box = shapes.add_textbox(
        Inches(2), Inches(4.5), Inches(12), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = "如有任何問題，歡迎提出討論"
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.LEFT

def add_summary_slide(prs, summary_data):
    """Add a summary slide with the exact layout shown in image 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # 添加白色背景
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 純白色背景
    background.line.fill.background()
    
    # 添加標題
    title_box = shapes.add_textbox(
        Inches(1), Inches(0.5),
        Inches(14), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = "競品分析總結"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 144, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 三個方框的佈局
    column_width = Inches(4.5)
    column_height = Inches(5.5)
    gap = Inches(0.7)
    
    # 計算三個方框的起始位置，確保居中
    total_width = column_width * 3 + gap * 2
    start_left = (prs.slide_width - total_width) / 2
    top = Inches(2.2)  # 位於標題下方
    
    # ===== 左方框：關鍵數據 =====
    # 1. 創建方框
    data_box = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # 使用矩形，不是圓角矩形
        start_left, top, column_width, column_height
    )
    data_box.fill.solid()
    data_box.fill.fore_color.rgb = RGBColor(240, 248, 255)  # 淺藍色背景
    data_box.line.color.rgb = RGBColor(173, 216, 230)  # 藍色邊框，較淺
    
    # 2. 添加標題
    data_title = shapes.add_textbox(
        start_left, top + Inches(0.4),
        column_width, Inches(0.6)
    )
    data_title_para = data_title.text_frame.add_paragraph()
    data_title_para.text = "關鍵數據"
    data_title_para.font.size = Pt(24)
    data_title_para.font.bold = True
    data_title_para.font.color.rgb = RGBColor(30, 144, 255)
    data_title_para.alignment = PP_ALIGN.CENTER
    
    # 4. 添加內容
    data_content = shapes.add_textbox(
        start_left + Inches(0.3), top + Inches(1.3),
        column_width - Inches(0.6), column_height - Inches(1.5)
    )
    data_content.text_frame.word_wrap = True
    
    for item in summary_data['dataSupport']:
        p = data_content.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(10)  # 段落間距
    
    # ===== 中間方框：關鍵發現 =====
    # 1. 創建方框
    middle_left = start_left + column_width + gap
    findings_box = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # 使用矩形
        middle_left, top, column_width, column_height
    )
    findings_box.fill.solid()
    findings_box.fill.fore_color.rgb = RGBColor(240, 255, 240)  # 淺綠色背景
    findings_box.line.color.rgb = RGBColor(144, 238, 144)  # 淺綠色邊框
    
    # 2. 添加標題
    findings_title = shapes.add_textbox(
        middle_left, top + Inches(0.4),
        column_width, Inches(0.6)
    )
    findings_title_para = findings_title.text_frame.add_paragraph()
    findings_title_para.text = "關鍵發現"
    findings_title_para.font.size = Pt(24)
    findings_title_para.font.bold = True
    findings_title_para.font.color.rgb = RGBColor(46, 204, 113)
    findings_title_para.alignment = PP_ALIGN.CENTER
    
    # 4. 添加內容
    findings_content = shapes.add_textbox(
        middle_left + Inches(0.3), top + Inches(1.3),
        column_width - Inches(0.6), column_height - Inches(1.5)
    )
    findings_content.text_frame.word_wrap = True
    
    for item in summary_data['keyFindings']:
        p = findings_content.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(10)  # 段落間距
    
    # ===== 右方框：具體建議 =====
    # 1. 創建方框
    right_left = middle_left + column_width + gap
    recom_box = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # 使用矩形
        right_left, top, column_width, column_height
    )
    recom_box.fill.solid()
    recom_box.fill.fore_color.rgb = RGBColor(255, 248, 240)  # 淺橙色背景
    recom_box.line.color.rgb = RGBColor(250, 214, 165)  # 淺橙色邊框
    
    # 2. 添加標題
    recom_title = shapes.add_textbox(
        right_left, top + Inches(0.4),
        column_width, Inches(0.6)
    )
    recom_title_para = recom_title.text_frame.add_paragraph()
    recom_title_para.text = "具體建議"
    recom_title_para.font.size = Pt(24)
    recom_title_para.font.bold = True
    recom_title_para.font.color.rgb = RGBColor(230, 126, 34)
    recom_title_para.alignment = PP_ALIGN.CENTER
    
    # 3. 添加內容
    recom_content = shapes.add_textbox(
        right_left + Inches(0.3), top + Inches(1.3),
        column_width - Inches(0.6), column_height - Inches(1.5)
    )
    recom_content.text_frame.word_wrap = True
    
    for item in summary_data['recommendations']:
        p = recom_content.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_after = Pt(10)  # 段落間距
    
    return slide

def add_app_header_slide(prs, app_name, app_logo_url):
    """創建應用程式分析的標題頁面，風格與總結分析頁面一致"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    
    # 添加淺藍灰色背景，與總結分析頁面一致
    background = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 250, 252)  # 淺藍灰色背景，與總結分析頁面相同
    background.line.fill.background()
    
    # 設定通用主色調
    main_color = RGBColor(30, 144, 255)  # 預設藍色
    logo_background_color = RGBColor(255, 255, 255)  # 預設白色背景
    
    # 嘗試下載和轉換圖片
    logo_path = None
    if app_logo_url:
        try:
            logo_path = download_and_convert_image(app_logo_url, app_name)
        except Exception as e:
            logger.error(f"Error downloading logo for app header slide: {str(e)}")
    
    # 垂直置中的位置計算
    vertical_center = prs.slide_height / 2
    logo_height = Inches(1.8)
    title_height = Inches(1.0)
    total_height = logo_height + Inches(0.5) + title_height  # logo + 間距 + 標題
    
    start_y = vertical_center - (total_height / 2)  # 從這個位置開始放置元素
    
    # 如果有Logo，放置在標題上方並垂直置中
    if logo_path and os.path.exists(logo_path):
        # 為Logo創建背景框（圓角矩形），水平置中
        logo_width = Inches(1.8)
        logo_left = (prs.slide_width - logo_width) / 2
        
        logo_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            logo_left, start_y,  # 水平置中，垂直位於計算的起始位置
            logo_width, logo_height  # 寬度和高度
        )
        logo_box.fill.solid()
        logo_box.fill.fore_color.rgb = logo_background_color
        logo_box.line.color.rgb = main_color
        
        # 添加應用標誌，在框內置中
        icon_width = Inches(1.4)
        icon_left = logo_left + (logo_width - icon_width) / 2
        icon_top = start_y + (logo_height - icon_width) / 2
        
        icon = shapes.add_picture(
            logo_path,
            icon_left, icon_top,  # 在框內水平和垂直置中
            width=icon_width  # 寬度
        )
        
        # 標題的位置向下移動到Logo下方
        title_top = start_y + logo_height + Inches(0.5)
    else:
        # 如果沒有Logo，標題直接垂直置中
        title_top = vertical_center - (title_height / 2)
    
    # 添加應用名稱標題 - 水平置中
    title_width = Inches(10)
    title_left = (prs.slide_width - title_width) / 2
    
    title_box = shapes.add_textbox(
        title_left, title_top,
        title_width, title_height
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    
    # 直接使用 app_name
    title_para.text = f"{app_name} 應用程式分析"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = main_color
    title_para.alignment = PP_ALIGN.CENTER  # 水平置中對齊
    
    # 清理臨時文件
    if logo_path and os.path.exists(logo_path):
        try:
            os.remove(logo_path)
            original_path = os.path.join(os.path.dirname(logo_path), f"{app_name}_original")
            if os.path.exists(original_path):
                os.remove(original_path)
            temp_dir = os.path.dirname(logo_path)
            if os.path.exists(temp_dir):
                os.rmdir(temp_dir)
        except Exception as cleanup_error:
            logger.warning(f"Error cleaning up temp files: {str(cleanup_error)}")
    
    return slide

def add_app_analysis_slide(prs, app):
    """向簡報添加應用分析頁面"""
    # 首先添加應用標題頁
    add_app_header_slide(prs, app['name'], app.get('logo', None))
    
    # 然後添加標準內容頁面
    # App Overview
    add_content_slide(
        prs,
        f"{app['name']} - 概述",
        {
            "基本資訊": {
                "iOS 評分": f"{app['ratings']['ios']} ⭐",
                "Android 評分": f"{app['ratings']['android']} ⭐",
                "總評論數": f"{app['reviews']['count']} 則"
            },
            "核心功能": app['features']['core'],
            "主要優勢": app['features']['advantages'],
            "待改進項目": app['features']['improvements']
        },
        app.get('logo', None),
        app['name']
    )
    
    # App UX Analysis
    add_content_slide(
        prs,
        f"{app['name']} - 用戶體驗分析",
        {
            "用戶體驗評分": {
                "會員登入": f"{app['uxScores']['memberlogin']}%",
                "搜尋功能": f"{app['uxScores']['search']}%",
                "商品相關": f"{app['uxScores']['product']}%",
                "結帳付款": f"{app['uxScores']['checkout']}%",
                "客戶服務": f"{app['uxScores']['service']}%",
                "其他功能": f"{app['uxScores']['other']}%"
            },
            "優勢分析": app['uxAnalysis']['strengths'],
            "改進建議": app['uxAnalysis']['improvements'],
            "分析摘要": app['uxAnalysis']['summary']
        },
        app.get('logo', None),
        app['name']
    )
    
    # App Review Analysis
    add_content_slide(
        prs,
        f"{app['name']} - 評論分析",
        {
            "評論統計": {
                "正面評價": f"{app['reviews']['stats']['positive']}%",
                "負面評價": f"{app['reviews']['stats']['negative']}%"
            },
            "用戶好評項目": app['reviews']['analysis']['advantages'],
            "用戶反饋問題": app['reviews']['analysis']['improvements'],
            "評論分析摘要": app['reviews']['analysis']['summary']
        },
        app.get('logo', None),
        app['name']
    )

def generate_competitive_analysis_ppt(input_file: str, output_file: str):
    """Generate a competitive analysis PowerPoint presentation."""
    try:
        logger.info(f"Starting PPT generation with input file: {input_file}")
        
        # Load and validate input data
        try:
            with open(input_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
                logger.info("Successfully loaded input JSON file")
        except Exception as e:
            logger.error(f"Error loading input file: {str(e)}")
            raise ValueError(f"Failed to load input file: {str(e)}")

        # Validate required fields
        required_fields = ['title', 'date', 'apps']
        missing_fields = [field for field in required_fields if field not in data]
        if missing_fields:
            error_msg = f"Missing required fields in input data: {', '.join(missing_fields)}"
            logger.error(error_msg)
            raise ValueError(error_msg)

        # Create presentation
        try:
            prs = Presentation()
            set_slide_size_to_16x9(prs)
            logger.info("Created new presentation with 16:9 aspect ratio")
        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            raise RuntimeError(f"Failed to create presentation: {str(e)}")

        # Add title slide
        try:
            add_title_slide(prs, data['title'], data['date'])
            logger.info("Added title slide")
        except Exception as e:
            logger.error(f"Error adding title slide: {str(e)}")
            raise RuntimeError(f"Failed to add title slide: {str(e)}")

        # Add content for each app
        try:
            for app in data['apps']:
                logger.info(f"Processing app: {app.get('name', 'Unknown')}")
                add_app_analysis_slide(prs, app)
        except Exception as e:
            logger.error(f"Error processing app data: {str(e)}")
            raise RuntimeError(f"Failed to process app data: {str(e)}")

        # Add summary slide
        try:
            if 'summary' in data:
                add_summary_slide(prs, data['summary'])
                logger.info("Added summary slide")
        except Exception as e:
            logger.error(f"Error adding summary slide: {str(e)}")
            raise RuntimeError(f"Failed to add summary slide: {str(e)}")

        # Add ending slide
        try:
            add_ending_slide(prs)
            logger.info("Added ending slide")
        except Exception as e:
            logger.error(f"Error adding ending slide: {str(e)}")
            raise RuntimeError(f"Failed to add ending slide: {str(e)}")

        # Save presentation
        try:
            # Ensure the output directory exists
            os.makedirs(os.path.dirname(output_file), exist_ok=True)
            prs.save(output_file)
            logger.info(f"Successfully saved presentation to: {output_file}")
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise RuntimeError(f"Failed to save presentation: {str(e)}")

    except Exception as e:
        logger.error(f"Error in generate_competitive_analysis_ppt: {str(e)}", exc_info=True)
        raise

    return output_file 